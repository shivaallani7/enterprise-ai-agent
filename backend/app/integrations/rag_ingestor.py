"""
RAG Document Ingestion Pipeline.

Supported types: PDF, Word (.docx), Excel / CSV / TSV, PowerPoint (.pptx),
                 Markdown, plain text (.txt / .log), images (Vision LLM).

Pipeline (per spec):
  1.  Identity resolution  — Cosmos DB registry → INSERT vs UPDATE
  2.  Content hash         — skip if file unchanged
  3.  Delete stale vectors — UPDATE mode only
  4.  Extract content      — per-type extractor
  5.  Clean & normalize    — unicode, whitespace, boilerplate
  6.  Chunk               — per-type strategy
  7.  Metadata enrichment — full metadata dict per chunk
  8.  PII detection       — regex patterns → redact
  9.  Exact dedup         — hash within batch
  10. Dense embedding     — batched via OpenAI / Azure OpenAI
  11. Upsert              — Azure AI Search docs-index
  12. Registry update     — Cosmos DB
  13. IngestReport        — returned to caller / API
"""
from __future__ import annotations

import base64
import hashlib
import io
import re
import time
import unicodedata
import uuid
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import structlog

from app.config import get_settings
from app.memory.cosmos_store import IngestRegistry

logger = structlog.get_logger()
settings = get_settings()

# ── Limits & defaults ──────────────────────────────────────────────────────────

MAX_FILE_BYTES  = 500 * 1024 * 1024   # 500 MB hard limit
WARN_FILE_BYTES = 100 * 1024 * 1024   # 100 MB soft warning

CHUNK_SIZE_TOKENS = 512
OVERLAP_PCT       = 0.20
MIN_CHUNK_TOKENS  = 10   # very low — orphan merge happens inside _semantic_chunk
ROWS_PER_CHUNK    = 15
EMBED_BATCH_SIZE  = 16
UPSERT_BATCH_SIZE = 100

_EXT_TO_DOCTYPE: dict[str, str] = {
    ".pdf":  "pdf",
    ".docx": "word",  ".doc":  "word",
    ".xlsx": "excel", ".xls":  "excel", ".csv": "excel", ".tsv": "excel",
    ".pptx": "pptx",  ".ppt":  "pptx",
    ".md":   "markdown", ".mdx": "markdown",
    ".txt":  "txt",   ".log":  "txt",
    ".png":  "image", ".jpg":  "image", ".jpeg": "image",
    ".gif":  "image", ".webp": "image", ".svg":  "image",
    ".html": "txt",   ".htm":  "txt",
}

_MIME_TO_EXT: dict[str, str] = {
    "application/pdf":                                                          ".pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document":  ".docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":        ".xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation":".pptx",
    "text/csv":        ".csv",
    "text/plain":      ".txt",
    "text/markdown":   ".md",
    "image/png":       ".png",
    "image/jpeg":      ".jpg",
    "image/gif":       ".gif",
    "image/webp":      ".webp",
}

# PII patterns: SSN, credit card (loose), passport-style
_PII_RE = [
    re.compile(r"\b\d{3}-\d{2}-\d{4}\b"),
    re.compile(r"\b(?:\d[ -]?){13,16}\b"),
    re.compile(r"\b[A-Z]{1,2}\d{6,9}\b"),
]

# ── Data classes ───────────────────────────────────────────────────────────────

@dataclass
class Chunk:
    chunk_id: str
    doc_id:   str
    text:     str
    metadata: dict[str, Any]


@dataclass
class IngestReport:
    doc_id:          str
    file_name:       str
    doc_type:        str
    mode:            str    # INSERT | UPDATE | SKIPPED
    status:          str    # SUCCESS | FAILED | SKIPPED | EMPTY_DOCUMENT | UNSUPPORTED_TYPE
    chunks_created:  int = 0
    chunks_skipped:  int = 0
    chunks_failed:   int = 0
    embedding_model: str = ""
    chunk_strategy:  str = ""
    chunk_size:      int = CHUNK_SIZE_TOKENS
    overlap_pct:     int = int(OVERLAP_PCT * 100)
    version:         int = 1
    processing_ms:   int = 0
    warnings:        list[str] = field(default_factory=list)
    indexed_at:      str = ""
    error:           str = ""


# ── Token helpers ──────────────────────────────────────────────────────────────

def _count_tokens(text: str) -> int:
    try:
        import tiktoken
        enc = tiktoken.get_encoding("cl100k_base")
        return len(enc.encode(text))
    except Exception:
        return len(text) // 4   # ~4 chars per token fallback


# ── Text cleaning (Step 5) ─────────────────────────────────────────────────────

def _clean(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\x00", "")
    text = unicodedata.normalize("NFC", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    # Normalize smart quotes
    text = text.replace("\u201c", '"').replace("\u201d", '"')
    text = text.replace("\u2018", "'").replace("\u2019", "'")
    # Strip common boilerplate
    text = re.sub(r"Page \d+ of \d+", "", text, flags=re.IGNORECASE)
    text = re.sub(r"Confidential\s*[\-–]\s*Internal only", "", text, flags=re.IGNORECASE)
    return text.strip()


# ── PII (Step 8) ───────────────────────────────────────────────────────────────

def _detect_pii(text: str) -> bool:
    return any(p.search(text) for p in _PII_RE)


def _redact_pii(text: str) -> str:
    for p in _PII_RE:
        text = p.sub("[REDACTED]", text)
    return text


# ── Language detection ─────────────────────────────────────────────────────────

def _detect_lang(text: str) -> str:
    try:
        from langdetect import detect
        return detect(text[:2000]) or "en"
    except Exception:
        return "en"


# ── Chunking strategies ────────────────────────────────────────────────────────

def _token_split(text: str, max_tok: int, overlap_tok: int) -> list[str]:
    """Split text into fixed-size token windows with overlap."""
    words = text.split(" ")
    chunks: list[str] = []
    i = 0
    while i < len(words):
        # Greedy fill up to max_tok (approximating 1 word ≈ 1 token)
        j = i
        tok = 0
        while j < len(words) and tok < max_tok:
            tok += 1
            j += 1
        chunk = " ".join(words[i:j]).strip()
        if chunk:
            chunks.append(chunk)
        # Step forward minus overlap
        step = max(j - i - overlap_tok, 1)
        i += step
    return chunks


def _semantic_chunk(text: str, section_title: str = "") -> list[dict]:
    """
    Semantic chunking: split at H1/H2/H3 headings first, then paragraphs,
    then recursively by tokens if still too large. Carries overlap tail.
    """
    heading_re = re.compile(r"^(#{1,3} .+)$", re.MULTILINE)
    parts = heading_re.split(text)
    # parts alternates: body, heading, body, heading, ...
    # Build (heading, body) pairs
    sections: list[tuple[str, str]] = []
    i = 0
    while i < len(parts):
        if heading_re.match(parts[i]):
            hdr = parts[i].lstrip("#").strip()
            body = parts[i + 1] if i + 1 < len(parts) else ""
            sections.append((hdr, body))
            i += 2
        else:
            sections.append((section_title, parts[i]))
            i += 1

    overlap_tok = int(CHUNK_SIZE_TOKENS * OVERLAP_PCT)
    chunks: list[dict] = []
    tail = ""

    for hdr, body in sections:
        paras = [p.strip() for p in re.split(r"\n\n+", body) if p.strip()]
        current = tail
        for para in paras:
            candidate = (current + "\n\n" + para).strip() if current else para
            if _count_tokens(candidate) <= CHUNK_SIZE_TOKENS:
                current = candidate
            else:
                if current and _count_tokens(current) >= MIN_CHUNK_TOKENS:
                    chunks.append({"text": current, "section_title": hdr})
                if _count_tokens(para) > CHUNK_SIZE_TOKENS:
                    for sub in _token_split(para, CHUNK_SIZE_TOKENS, overlap_tok):
                        if _count_tokens(sub) >= MIN_CHUNK_TOKENS:
                            chunks.append({"text": sub, "section_title": hdr})
                    current = sub if _token_split(para, CHUNK_SIZE_TOKENS, overlap_tok) else ""
                else:
                    current = para

        if current:
            if chunks and _count_tokens(current) < MIN_CHUNK_TOKENS:
                # Merge tiny trailing fragment into previous chunk
                chunks[-1]["text"] = chunks[-1]["text"] + "\n" + current
            else:
                chunks.append({"text": current, "section_title": hdr})
            words = current.split()
            tail = " ".join(words[-overlap_tok:]) if len(words) > overlap_tok else current
        else:
            tail = ""

    return chunks


def _row_chunk(rows: list[str], headers: str, sheet_name: str) -> list[dict]:
    """Group spreadsheet rows into batches, always prepend headers."""
    chunks: list[dict] = []
    for i in range(0, len(rows), ROWS_PER_CHUNK):
        batch = rows[i: i + ROWS_PER_CHUNK]
        text = f"Columns: {headers}\n" + "\n".join(batch)
        chunks.append({
            "text":           text,
            "section_title":  sheet_name,
            "sheet_name":     sheet_name,
            "row_range":      f"{i + 1}-{i + len(batch)}",
            "content_type":   "table",
            "chunk_strategy": "row_batch",
        })
    return chunks


# ── Extractors ─────────────────────────────────────────────────────────────────

async def _extract_pdf(data: bytes) -> list[dict]:
    try:
        import pdfplumber
    except ImportError:
        return [{"text": "[pdfplumber not installed — pip install pdfplumber]", "page": 1}]

    pages: list[dict] = []
    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            total_chars = 0
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                # Inline table extraction
                for tbl in (page.extract_tables() or []):
                    if not tbl:
                        continue
                    hdrs = " | ".join(str(c or "") for c in (tbl[0] or []))
                    for row in tbl[1:]:
                        text += "\n" + " | ".join(str(c or "") for c in row)
                total_chars += len(text)
                if text.strip():
                    pages.append({"text": text, "page": i})

            avg = total_chars / max(len(pdf.pages), 1)
            if avg < 50 and pages:
                pages[0]["text"] = (
                    f"[SCANNED PDF — avg {avg:.0f} chars/page, OCR required]\n"
                    + pages[0]["text"]
                )
    except Exception as exc:
        pages.append({"text": f"[PDF extraction error: {exc}]", "page": 1})
    return pages


def _extract_docx(data: bytes) -> list[dict]:
    try:
        from docx import Document
        from docx.oxml.ns import qn
    except ImportError:
        return [{"text": "[python-docx not installed — pip install python-docx]"}]

    try:
        doc = Document(io.BytesIO(data))
    except Exception as exc:
        # Common for .doc (old binary format) — fall back to raw XML text
        return [{"text": f"[Could not open as .docx: {exc}. If this is a .doc file, save it as .docx first.]"}]

    sections: list[dict] = []
    current_text = ""
    current_title = ""

    for para in doc.paragraphs:
        style = para.style.name if para.style else ""
        # Collect runs directly — para.text can miss some formatted runs
        text = "".join(run.text for run in para.runs).strip() or para.text.strip()
        if not text:
            continue
        if style.startswith("Heading") or style.startswith("Title"):
            if current_text:
                sections.append({"text": current_text.strip(), "section_title": current_title})
            current_title = text
            current_text = f"[Section: {text}]\n"
        else:
            current_text += text + "\n"

    if current_text:
        sections.append({"text": current_text.strip(), "section_title": current_title})

    # Extract text from text boxes (shapes / drawing objects)
    try:
        for shape_text in _extract_docx_textboxes(doc):
            if shape_text.strip():
                sections.append({"text": shape_text.strip(), "section_title": "Text Box"})
    except Exception:
        pass

    # Tables
    for tbl in doc.tables:
        if not tbl.rows:
            continue
        hdrs = " | ".join(c.text.strip() for c in tbl.rows[0].cells)
        rows = [" | ".join(c.text.strip() for c in row.cells) for row in tbl.rows[1:]]
        # Include header row too as standalone if body is empty
        if rows:
            sections.append({
                "text":         f"Table:\n{hdrs}\n" + "\n".join(rows),
                "section_title": "Table",
                "content_type": "table",
            })
        elif hdrs.strip():
            sections.append({"text": f"Table header: {hdrs}", "section_title": "Table"})

    # Last resort: if nothing extracted, pull raw XML text
    if not sections or all(not s.get("text", "").strip() for s in sections):
        try:
            raw = "\n".join(
                elem.text or ""
                for elem in doc.element.iter()
                if elem.text and elem.text.strip()
            )
            if raw.strip():
                sections = [{"text": raw.strip(), "section_title": ""}]
        except Exception:
            pass

    return sections if sections else [{"text": "", "section_title": ""}]


def _extract_docx_textboxes(doc) -> list[str]:
    """Extract text from Word text boxes (drawing/txbx elements)."""
    from lxml import etree
    texts: list[str] = []
    ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    draw_ns = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
    for body_elem in doc.element.body:
        for txbx in body_elem.iter("{http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing}txbxContent"):
            parts = []
            for t in txbx.iter(f"{{{ns}}}t"):
                if t.text:
                    parts.append(t.text)
            if parts:
                texts.append(" ".join(parts))
    return texts


def _extract_excel(data: bytes, filename: str) -> list[dict]:
    import pandas as pd

    ext = Path(filename).suffix.lower()
    sheets: list[dict] = []
    try:
        if ext in (".csv", ".tsv"):
            sep = "\t" if ext == ".tsv" else ","
            df = pd.read_csv(io.BytesIO(data), sep=sep, dtype=str, na_filter=False)
            df = df.dropna(how="all").dropna(axis=1, how="all")
            sheets.append({"name": "Sheet1", "df": df})
        else:
            xl = pd.read_excel(io.BytesIO(data), sheet_name=None, dtype=str, na_filter=False)
            for sname, df in xl.items():
                df = df.dropna(how="all").dropna(axis=1, how="all")
                sheets.append({"name": sname, "df": df})
    except Exception as exc:
        return [{"text": f"[Excel extraction error: {exc}]"}]

    chunks: list[dict] = []
    for sheet in sheets:
        df = sheet["df"]
        sname = sheet["name"]
        if df.empty:
            continue
        hdrs = " | ".join(str(c) for c in df.columns)
        rows = [
            f"Sheet: {sname} | " + " | ".join(
                f"{col}: {val}" for col, val in row.items() if str(val).strip()
            )
            for _, row in df.iterrows()
        ]
        chunks.extend(_row_chunk(rows, hdrs, sname))
    return chunks


def _extract_pptx(data: bytes) -> list[dict]:
    try:
        from pptx import Presentation
    except ImportError:
        return [{"text": "[python-pptx not installed — pip install python-pptx]", "slide": 1}]

    prs = Presentation(io.BytesIO(data))
    slides: list[dict] = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        body_parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            t = shape.text_frame.text.strip()
            if not t:
                continue
            if "title" in shape.name.lower() and not title:
                title = t
            else:
                body_parts.append(t)
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        text = f"Slide {i}: {title}\n" + "\n".join(body_parts)
        if notes:
            text += f"\n[Speaker notes: {notes}]"
        if text.strip():
            slides.append({
                "text":           text,
                "section_title":  title or f"Slide {i}",
                "chunk_strategy": "slide_level",
            })
    return slides


async def _extract_image(data: bytes, filename: str, openai_client) -> str:
    if openai_client is None:
        return f"[Image: {filename} — vision model not available]"

    ext = Path(filename).suffix.lower().lstrip(".")
    mime_map = {"jpg": "image/jpeg", "jpeg": "image/jpeg",
                "png": "image/png", "gif": "image/gif", "webp": "image/webp"}
    mime = mime_map.get(ext, "image/png")
    b64 = base64.b64encode(data).decode()

    prompt = (
        "You are analyzing a technical diagram or flowchart for a RAG knowledge base. "
        "Provide a detailed description covering: "
        "1. The overall purpose of this diagram. "
        "2. All nodes, boxes, or components — their labels and roles. "
        "3. All connections, arrows, or edges and what they represent. "
        "4. Any decision points, conditions, or branching logic. "
        "5. Any annotations, legends, or callout labels. "
        "6. The overall flow direction (top-down, left-right, circular, etc.). "
        "Be precise and technical. Output only the description, no preamble."
    )
    model = settings.openai_model if settings.llm_provider == "openai" else settings.azure_openai_deployment
    try:
        resp = await openai_client.chat.completions.create(
            model=model,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
                ],
            }],
            max_tokens=600,
        )
        caption = resp.choices[0].message.content or ""
    except Exception as exc:
        caption = f"[Vision captioning failed: {exc}]"

    return f"[DIAGRAM DESCRIPTION]\n{caption}\n\n[IMAGE FILE]\n{filename}"


# ── Main pipeline ──────────────────────────────────────────────────────────────

class RagIngestor:
    """
    Full RAG ingestion pipeline for user-uploaded documents.

    Usage:
        ingestor = RagIngestor()
        report = await ingestor.ingest(file_bytes, filename)
    """

    def __init__(self) -> None:
        self._registry: IngestRegistry | None = None
        self._openai = None

    # ── Lazy singletons ───────────────────────────────────────────────────────

    def _get_registry(self) -> IngestRegistry:
        if self._registry is None:
            self._registry = IngestRegistry()
        return self._registry

    def _get_openai(self):
        if self._openai is None:
            from openai import AsyncOpenAI, AsyncAzureOpenAI
            if settings.llm_provider == "openai":
                self._openai = AsyncOpenAI(api_key=settings.openai_api_key)
            else:
                self._openai = AsyncAzureOpenAI(
                    azure_endpoint=settings.azure_openai_endpoint,
                    api_key=settings.azure_openai_api_key,
                    api_version=settings.azure_openai_api_version,
                )
        return self._openai

    # ── Public entry point ────────────────────────────────────────────────────

    async def ingest(self, file_bytes: bytes, filename: str) -> IngestReport:
        start_ms = int(time.time() * 1000)
        warnings: list[str] = []

        # Guard: file size
        if len(file_bytes) > MAX_FILE_BYTES:
            return IngestReport(
                doc_id="", file_name=filename, doc_type="", mode="INSERT",
                status="FAILED", error="File exceeds 500 MB limit",
                indexed_at=_iso_now(), processing_ms=_elapsed(start_ms),
            )
        if len(file_bytes) > WARN_FILE_BYTES:
            warnings.append(f"Large file ({len(file_bytes)//1024//1024} MB) — ingestion may be slow")

        # Detect doc type
        ext = Path(filename).suffix.lower()
        doc_type = _EXT_TO_DOCTYPE.get(ext, "txt")

        # ── Step 1: Identity resolution ───────────────────────────────────────
        registry = self._get_registry()
        record = await registry.lookup(filename)
        if record:
            mode = "UPDATE"
            doc_id = record["doc_id"]
            version = record.get("version", 1) + 1
        else:
            mode = "INSERT"
            slug = re.sub(r"[^a-z0-9]+", "-", filename.lower())[:30].strip("-")
            doc_id = f"upload-{slug}-{uuid.uuid4().hex[:6]}"
            version = 1

        # ── Step 2: Content hash ──────────────────────────────────────────────
        file_hash = hashlib.sha256(file_bytes).hexdigest()
        if mode == "UPDATE" and record.get("file_hash") == file_hash:
            return IngestReport(
                doc_id=doc_id, file_name=filename, doc_type=doc_type,
                mode=mode, status="SKIPPED",
                version=record.get("version", 1),
                processing_ms=_elapsed(start_ms), indexed_at=_iso_now(),
            )

        # ── Step 3: Delete stale vectors ──────────────────────────────────────
        if mode == "UPDATE":
            try:
                await self._delete_by_doc_id(doc_id)
                logger.info("Deleted stale vectors", doc_id=doc_id)
            except Exception as exc:
                warnings.append(f"Stale vector deletion incomplete: {exc}")

        # ── Step 4: Extract ───────────────────────────────────────────────────
        try:
            raw_chunks = await self._extract(file_bytes, filename, doc_type)
        except Exception as exc:
            return IngestReport(
                doc_id=doc_id, file_name=filename, doc_type=doc_type,
                mode=mode, status="FAILED", error=f"Extraction failed: {exc}",
                indexed_at=_iso_now(), processing_ms=_elapsed(start_ms),
            )

        if not raw_chunks:
            return IngestReport(
                doc_id=doc_id, file_name=filename, doc_type=doc_type,
                mode=mode, status="EMPTY_DOCUMENT",
                processing_ms=_elapsed(start_ms), indexed_at=_iso_now(),
            )

        # ── Step 5-7: Clean + enrich metadata ─────────────────────────────────
        lang = _detect_lang(raw_chunks[0].get("text", "") if raw_chunks else "")
        total = len(raw_chunks)
        chunks: list[Chunk] = []

        for i, rc in enumerate(raw_chunks):
            text = _clean(rc.get("text", ""))
            if not text:
                continue
            if _detect_pii(text):
                text = _redact_pii(text)
                warnings.append(f"PII detected and redacted in chunk {i + 1}")

            chunk_id = f"{doc_id}__chunk_{i:04d}"
            meta: dict[str, Any] = {
                "doc_id":            doc_id,
                "chunk_id":          chunk_id,
                "source":            "upload",
                "doc_type":          doc_type,
                "file_name":         filename,
                "file_hash":         f"sha256:{file_hash[:16]}",
                "section_title":     rc.get("section_title", ""),
                "page_number":       rc.get("page", 0),
                "sheet_name":        rc.get("sheet_name", ""),
                "row_range":         rc.get("row_range", ""),
                "language":          lang,
                "content_type":      rc.get("content_type", "prose"),
                "extraction_method": rc.get("extraction_method", "native"),
                "chunk_index":       i,
                "total_chunks":      total,
                "chunk_strategy":    rc.get("chunk_strategy", "semantic"),
                "version":           version,
            }
            chunks.append(Chunk(chunk_id=chunk_id, doc_id=doc_id, text=text, metadata=meta))

        if not chunks:
            return IngestReport(
                doc_id=doc_id, file_name=filename, doc_type=doc_type,
                mode=mode, status="EMPTY_DOCUMENT",
                processing_ms=_elapsed(start_ms), indexed_at=_iso_now(),
            )

        # ── Step 8 (exact dedup) ──────────────────────────────────────────────
        seen: set[str] = set()
        deduped: list[Chunk] = []
        skipped = 0
        for ch in chunks:
            h = hashlib.md5(ch.text.encode()).hexdigest()
            if h in seen:
                skipped += 1
            else:
                seen.add(h)
                deduped.append(ch)
        chunks = deduped

        # ── Step 9: Embed ─────────────────────────────────────────────────────
        embedding_model = settings.azure_openai_embedding_deployment or "none"
        if settings.azure_openai_embedding_deployment:
            await self._embed(chunks, warnings)

        # ── Step 10: Upsert ───────────────────────────────────────────────────
        failed = 0
        try:
            failed = await self._upsert(chunks)
        except Exception as exc:
            return IngestReport(
                doc_id=doc_id, file_name=filename, doc_type=doc_type,
                mode=mode, status="FAILED", error=f"Upsert failed: {exc}",
                chunks_created=len(chunks) - failed, chunks_failed=failed,
                warnings=warnings, processing_ms=_elapsed(start_ms), indexed_at=_iso_now(),
            )

        # ── Step 11: Registry update ──────────────────────────────────────────
        await registry.upsert({
            "id":          filename,
            "doc_id":      doc_id,
            "file_name":   filename,
            "file_hash":   file_hash,
            "doc_type":    doc_type,
            "chunk_count": len(chunks),
            "version":     version,
            "status":      "active",
            "indexed_at":  _iso_now(),
        })

        strategy = raw_chunks[0].get("chunk_strategy", "semantic") if raw_chunks else "semantic"
        return IngestReport(
            doc_id=doc_id,
            file_name=filename,
            doc_type=doc_type,
            mode=mode,
            status="SUCCESS",
            chunks_created=len(chunks) - failed,
            chunks_skipped=skipped,
            chunks_failed=failed,
            embedding_model=embedding_model,
            chunk_strategy=strategy,
            chunk_size=CHUNK_SIZE_TOKENS,
            overlap_pct=int(OVERLAP_PCT * 100),
            version=version,
            processing_ms=_elapsed(start_ms),
            warnings=warnings,
            indexed_at=_iso_now(),
        )

    # ── Private: extract ──────────────────────────────────────────────────────

    async def _extract(self, data: bytes, filename: str, doc_type: str) -> list[dict]:
        if doc_type == "pdf":
            pages = await _extract_pdf(data)
            chunks: list[dict] = []
            for p in pages:
                for sub in _semantic_chunk(p["text"]):
                    sub["page"] = p.get("page", 0)
                    sub["chunk_strategy"] = "semantic"
                    chunks.append(sub)
            return chunks

        if doc_type == "word":
            sections = _extract_docx(data)
            chunks = []
            for sec in sections:
                if sec.get("content_type") == "table":
                    sec.setdefault("chunk_strategy", "row_batch")
                    chunks.append(sec)
                else:
                    raw_text = sec.get("text", "").strip()
                    if not raw_text:
                        continue
                    sub_chunks = _semantic_chunk(raw_text, sec.get("section_title", ""))
                    if sub_chunks:
                        for sub in sub_chunks:
                            sub["chunk_strategy"] = "semantic"
                            chunks.append(sub)
                    else:
                        # Section text is non-empty but below chunking threshold — keep as-is
                        chunks.append({
                            "text":           raw_text,
                            "section_title":  sec.get("section_title", ""),
                            "chunk_strategy": "semantic",
                        })
            return chunks

        if doc_type == "excel":
            return _extract_excel(data, filename)

        if doc_type == "pptx":
            return [s for s in _extract_pptx(data) if s.get("text", "").strip()]

        if doc_type == "markdown":
            text = data.decode("utf-8", errors="replace")
            chunks = []
            for sub in _semantic_chunk(text):
                sub["chunk_strategy"] = "semantic"
                chunks.append(sub)
            return chunks

        if doc_type == "image":
            openai = self._get_openai()
            caption = await _extract_image(data, filename, openai)
            return [{
                "text":             caption,
                "section_title":    filename,
                "content_type":     "caption",
                "extraction_method":"vision_llm",
                "chunk_strategy":   "caption_as_chunk",
            }]

        # txt / html / unknown — fixed-token fallback
        text = data.decode("utf-8", errors="replace")
        overlap_tok = int(CHUNK_SIZE_TOKENS * 0.15)
        parts = _token_split(text, CHUNK_SIZE_TOKENS, overlap_tok)
        return [{"text": p, "chunk_strategy": "fixed_token"} for p in parts if p.strip()]

    # ── Private: embed ────────────────────────────────────────────────────────

    async def _embed(self, chunks: list[Chunk], warnings: list[str]) -> None:
        from tenacity import (
            AsyncRetrying, retry_if_exception_type,
            stop_after_attempt, wait_exponential_jitter,
        )
        from openai import RateLimitError, APIConnectionError

        openai = self._get_openai()
        deployment = settings.azure_openai_embedding_deployment

        for i in range(0, len(chunks), EMBED_BATCH_SIZE):
            batch = chunks[i: i + EMBED_BATCH_SIZE]
            texts = [f"Represent this document for retrieval: {ch.text[:8000]}" for ch in batch]
            try:
                async for attempt in AsyncRetrying(
                    retry=retry_if_exception_type((RateLimitError, APIConnectionError)),
                    stop=stop_after_attempt(3),
                    wait=wait_exponential_jitter(initial=1, max=30),
                    reraise=True,
                ):
                    with attempt:
                        resp = await openai.embeddings.create(input=texts, model=deployment)
                for ch, item in zip(batch, resp.data):
                    ch.metadata["content_vector"] = item.embedding
            except Exception as exc:
                warnings.append(f"Embedding batch {i // EMBED_BATCH_SIZE + 1} failed: {exc}")

    # ── Private: upsert ───────────────────────────────────────────────────────

    async def _upsert(self, chunks: list[Chunk]) -> int:
        from azure.core.credentials import AzureKeyCredential
        from azure.search.documents.aio import SearchClient as AzureSearchClient

        credential = AzureKeyCredential(settings.azure_search_api_key)
        failed = 0

        for i in range(0, len(chunks), UPSERT_BATCH_SIZE):
            batch = chunks[i: i + UPSERT_BATCH_SIZE]
            docs = []
            for ch in batch:
                m = ch.metadata
                doc: dict[str, Any] = {
                    "id":               ch.chunk_id,
                    "filepath":         m.get("file_name", ""),
                    "title":            m.get("section_title", "") or m.get("file_name", ""),
                    "content":          ch.text,
                    "chunk_type":       m.get("content_type", "prose"),
                    "source":           "upload",
                    "doc_id":           m.get("doc_id", ""),
                    "doc_type":         m.get("doc_type", ""),
                    "file_name":        m.get("file_name", ""),
                    "file_hash":        m.get("file_hash", ""),
                    "section_title":    m.get("section_title", ""),
                    "page_number":      m.get("page_number", 0),
                    "sheet_name":       m.get("sheet_name", ""),
                    "language":         m.get("language", "en"),
                    "extraction_method":m.get("extraction_method", "native"),
                    "chunk_index":      m.get("chunk_index", 0),
                    "total_chunks":     m.get("total_chunks", 1),
                    "chunk_strategy":   m.get("chunk_strategy", "semantic"),
                    "version":          m.get("version", 1),
                }
                if "content_vector" in m:
                    doc["content_vector"] = m["content_vector"]
                docs.append(doc)

            try:
                async with AzureSearchClient(
                    endpoint=settings.azure_search_endpoint,
                    index_name=settings.azure_search_docs_index,
                    credential=credential,
                ) as client:
                    results = await client.upload_documents(documents=docs)
                    batch_failed = sum(1 for r in results if not r.succeeded)
                    if batch_failed:
                        logger.warning("Some docs failed to upsert", failed=batch_failed, batch_start=i)
                    failed += batch_failed
            except Exception as exc:
                logger.error("Upsert batch error", error=str(exc), batch_start=i)
                failed += len(batch)

        return failed

    # ── Private: delete by doc_id ─────────────────────────────────────────────

    async def _delete_by_doc_id(self, doc_id: str) -> None:
        from azure.core.credentials import AzureKeyCredential
        from azure.search.documents.aio import SearchClient as AzureSearchClient

        credential = AzureKeyCredential(settings.azure_search_api_key)
        escaped = doc_id.replace("'", "''")
        async with AzureSearchClient(
            endpoint=settings.azure_search_endpoint,
            index_name=settings.azure_search_docs_index,
            credential=credential,
        ) as client:
            results = await client.search(
                search_text="*",
                filter=f"doc_id eq '{escaped}'",
                select=["id"],
                top=1000,
            )
            ids_to_delete = [{"id": r["id"]} async for r in results]
            if ids_to_delete:
                await client.delete_documents(documents=ids_to_delete)
                logger.info("Deleted stale chunks", count=len(ids_to_delete), doc_id=doc_id)


# ── Utilities ──────────────────────────────────────────────────────────────────

def _iso_now() -> str:
    return datetime.now(timezone.utc).isoformat()


def _elapsed(start_ms: int) -> int:
    return int(time.time() * 1000) - start_ms
